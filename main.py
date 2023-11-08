# coding: UTF-8

from pptx import Presentation
from pptx.util import Inches
import os
from PIL import Image
from pptx.util import Inches, Pt

# 画像のアスペクト比を保ちつつ、指定された幅にリサイズする関数
def resize_image(image_path, max_width):
    with Image.open(image_path) as img:
        # 元の画像のアスペクト比を保って最大幅に合わせてリサイズ
        original_width, original_height = img.size
        aspect_ratio = original_height / original_width
        new_width = max_width
        new_height = int(new_width * aspect_ratio)
        return new_width, new_height

presentation = Presentation()

# スライドのサイズを取得（インチ単位）
slide_width = presentation.slide_width
slide_height = presentation.slide_height

# 画像ファイルのリスト
image_paths = sorted([img for img in os.listdir('.') if img.endswith(".jpg")])

# 画像間のスペース（インチ単位）
image_margin = Inches(0.1)

# 1行あたりの画像数
images_per_row = 3

# 新しいスライドを追加する関数
def create_slide(presentation):
    blank_slide_layout = presentation.slide_layouts[5]  # 空白のスライドレイアウト

    return presentation.slides.add_slide(blank_slide_layout)

# スライドを追加
slide = create_slide(presentation)

# 画像を均等に配置するための幅を計算
available_width = slide_width - image_margin * (images_per_row + 1)
max_image_width = available_width / images_per_row

# 現在の行と列
current_row = 0
current_column = 0
left = image_margin

title_shape = slide.shapes.title
title_shape.text = "タイトル"
title_height = Inches(2.0)
top = top = title_height + image_margin

for image_path in image_paths:
    new_image_width, new_image_height = resize_image(image_path, max_image_width)

    slide.shapes.add_picture(image_path, left, top, new_image_width, new_image_height)

    # 次の画像の位置を更新
    left += new_image_width + image_margin
    current_column += 1

    if current_column >= images_per_row:
        current_column = 0
        left = image_margin
        top += new_image_height + image_margin
        current_row += 1

    # スライドの最大行数に達したら新しいスライドを追加
    if (top + new_image_height) > slide_height:
        slide = create_slide(presentation)
        current_row = 0
        current_column = 0
        left = image_margin
        top = image_margin


presentation.save('output.pptx')
